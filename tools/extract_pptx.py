#!/usr/bin/env python3
"""Extract text from .pptx files into per-slide markdown.

Usage:
  python3 tools/extract_pptx.py [paths...]
  (default: glob ../raw/*.pptx and ../*.pptx — auto-detects the source location)
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REPO = Path(__file__).resolve().parents[1]
EXTRACTED = REPO / "raw" / "extracted"

# Map source filename (loose) → output slug
SLUG_RULES = [
    (re.compile(r"chapter\s*9\b", re.I), "chapter-09-meiosis"),
    (re.compile(r"chapter\s*10\b", re.I), "chapter-10-inheritance"),
    (re.compile(r"chapter\s*11\b", re.I), "chapter-11-dna-structure-replication"),
    (re.compile(r"chapter\s*12\b", re.I), "chapter-12-how-genes-work"),
    (re.compile(r"chapter\s*13\b", re.I), "chapter-13-new-biology"),
]


def slug_for(name: str) -> str:
    for pat, slug in SLUG_RULES:
        if pat.search(name):
            return slug
    base = re.sub(r"\.pptx$", "", name, flags=re.I)
    return re.sub(r"[^a-z0-9-]+", "-", base.lower()).strip("-")


def extract_shape_text(shape, depth: int = 0) -> list[str]:
    """Recurse into shape, return list of text lines (stripped)."""
    out: list[str] = []
    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text and para.text:
                text = para.text.strip()
            if text:
                out.append(text)
    if shape.shape_type == 6 and hasattr(shape, "shapes"):  # group
        for sub in shape.shapes:
            out.extend(extract_shape_text(sub, depth + 1))
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(cells)
        if rows:
            out.append(format_table(rows))
    return out


def format_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join(["---"] * len(rows[0])) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(c.replace("\n", " ").replace("|", "\\|") for c in r) + " |")
    return "\n".join(out)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    out: list[str] = []
    out.append(f"# {path.name}\n")
    out.append(f"_Extracted from `{path.name}` — {len(prs.slides)} slides_\n")
    for i, slide in enumerate(prs.slides, 1):
        anchor = f"slide-{i:02d}"
        out.append(f"## Slide {i} {{#{anchor}}}\n")
        # Title (if any)
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        if title:
            out.append(f"**{title}**\n")
        # Other shape text (skip title shape to avoid duplication)
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            for line in extract_shape_text(shape):
                if line and line != title:
                    body_lines.append(line)
        for line in body_lines:
            if "\n" in line:  # already-formatted table
                out.append(line + "\n")
            else:
                out.append(f"- {line}")
        # Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.notes_slide.notes_text_frame else ""
            if notes:
                out.append("")
                out.append("### Notes\n")
                for para in notes.split("\n"):
                    p = para.strip()
                    if p:
                        out.append(f"> {p}")
        out.append("")
    return "\n".join(out)


def main() -> int:
    paths = [Path(p) for p in sys.argv[1:]]
    if not paths:
        # Auto-detect: look for *.pptx in REPO and REPO/raw and REPO parent (Documents/Study)
        for d in [REPO, REPO / "raw"]:
            if d.exists():
                paths.extend(sorted(d.glob("*.pptx")))
        paths = list({p.resolve() for p in paths})

    if not paths:
        print("no .pptx files found", file=sys.stderr)
        return 1

    EXTRACTED.mkdir(parents=True, exist_ok=True)
    for p in paths:
        slug = slug_for(p.name)
        out_path = EXTRACTED / f"{slug}.md"
        print(f"  {p.name} → {out_path.relative_to(REPO)}")
        out_path.write_text(extract_pptx(p), encoding="utf-8")
    print(f"extracted {len(paths)} pptx files into {EXTRACTED.relative_to(REPO)}/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
